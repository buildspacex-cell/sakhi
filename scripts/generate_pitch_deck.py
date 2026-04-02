"""
Generates Sakhi investor pitch deck (PowerPoint) from the /pitch video story.
Run: python3 scripts/generate_pitch_deck.py
Output: docs/Sakhi_Pitch_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x02, 0x06, 0x17)   # #020617 — main background
BG_CARD   = RGBColor(0x0b, 0x11, 0x20)   # slightly lighter card bg
ACCENT    = RGBColor(0x8c, 0xb7, 0xff)   # #8cb7ff — blue accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_DIM = RGBColor(0xB0, 0xBE, 0xD8)   # muted body text
WHITE_SUB = RGBColor(0x6A, 0x7A, 0x96)   # dimmer labels
GOLD      = RGBColor(0xFF, 0xD5, 0x9A)   # amber — ideas

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl

def txb(sl, text, left, top, width, height,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True, line_spacing=None):
    tf = sl.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as PPt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return tf

def rect(sl, left, top, width, height, fill_color=None, line_color=None, line_width_pt=1):
    shape = sl.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def pill(sl, text, left, top, width, height,
         fill=RGBColor(0x10,0x18,0x2e), border=RGBColor(0x30,0x45,0x70),
         text_color=WHITE_DIM, size=13):
    r = rect(sl, left, top, width, height, fill_color=fill, line_color=border, line_width_pt=0.5)
    tf = sl.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = text_color
    return r

def divider(sl, top, color=RGBColor(0x20,0x30,0x50)):
    rect(sl, Inches(1), top, SLIDE_W - Inches(2), Pt(1), fill_color=color)

def label_kicker(sl, text, left, top, width=Inches(8)):
    txb(sl, text, left, top, width, Inches(0.3),
        size=9, bold=True, color=ACCENT)

def multi_para(sl, lines, left, top, width, height,
               size=16, color=WHITE_DIM, leading=1.5, bold=False):
    """Add multiple lines as separate runs in one textbox."""
    box = sl.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

# ── SLIDE 1 — Hero ────────────────────────────────────────────────────────────
sl = add_slide()

# Central orb suggestion — concentric circles via shapes
cx, cy = SLIDE_W / 2, SLIDE_H / 2 - Inches(0.8)
for d, alpha_hint in [(Inches(3.2), RGBColor(0x10,0x18,0x30)),
                       (Inches(2.4), RGBColor(0x14,0x20,0x3e)),
                       (Inches(1.6), RGBColor(0x1a,0x2a,0x52))]:
    r = d / 2
    sh = sl.shapes.add_shape(9, cx - r, cy - r, d, d)  # 9 = oval
    sh.fill.solid()
    sh.fill.fore_color.rgb = alpha_hint
    sh.line.color.rgb = ACCENT
    sh.line.width = Pt(0.4 if d == Inches(3.2) else 0.6 if d == Inches(2.4) else 1.0)

# SAKHI word in orb
txb(sl, "SAKHI", cx - Inches(0.8), cy - Inches(0.22), Inches(1.6), Inches(0.44),
    size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Tagline below orb
txb(sl, "Your life, connected across time.",
    Inches(2), cy + Inches(1.1), Inches(9.33), Inches(0.6),
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(sl, "A companion intelligence that holds your life together so you can shape it.",
    Inches(2.5), cy + Inches(1.85), Inches(8.33), Inches(0.5),
    size=15, color=WHITE_DIM, align=PP_ALIGN.CENTER)

# Bottom — founders email
txb(sl, "founders@sakhiintelligence.com  ·  Sakhi · 2026",
    Inches(1), SLIDE_H - Inches(0.55), Inches(11.33), Inches(0.4),
    size=10, color=WHITE_SUB, align=PP_ALIGN.CENTER)

# ── SLIDE 2 — The Problem: The Noise ─────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "THE PROBLEM", Inches(1), Inches(0.7))

txb(sl, "There's a conversation\nhappening in your head.",
    Inches(1), Inches(1.1), Inches(7.5), Inches(2.2),
    size=46, bold=True, color=WHITE)

txb(sl, "It never really stops.",
    Inches(1), Inches(3.5), Inches(7), Inches(0.45), size=22, color=WHITE_DIM)
txb(sl, "And it doesn't carry forward.",
    Inches(1), Inches(4.0), Inches(7), Inches(0.45), size=22, color=WHITE_DIM)
txb(sl, "Nothing actually holds it together.",
    Inches(1), Inches(4.5), Inches(7), Inches(0.45), size=22, bold=True, color=WHITE)

# Right side — floating thought chips
thoughts = [
    ("did I reply?",        Inches(8.5), Inches(1.0)),
    ("✦ I just had a breakthrough", Inches(8.0), Inches(1.7)),
    ("flu shot for the kids", Inches(8.8), Inches(2.4)),
    ("this thread isn't finished", Inches(8.2), Inches(3.1)),
    ("pick up laundry",      Inches(9.0), Inches(3.8)),
    ("what did I promise?",  Inches(8.3), Inches(4.5)),
    ("✦ this connects it all", Inches(8.7), Inches(5.2)),
    ("need to follow up",    Inches(8.1), Inches(5.9)),
]
for t, x, y in thoughts:
    is_idea = t.startswith("✦")
    c = GOLD if is_idea else WHITE_DIM
    pill(sl, t, x, y, Inches(3.8), Inches(0.38),
         fill=RGBColor(0x12,0x1c,0x38) if is_idea else RGBColor(0x0e,0x14,0x26),
         border=RGBColor(0x50,0x60,0x90) if is_idea else RGBColor(0x22,0x30,0x50),
         text_color=c, size=12)

# ── SLIDE 3 — Breakdown ───────────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "THE BREAKDOWN", Inches(1), Inches(0.7))

txb(sl, "You wind down for the day.\nThe signal fades with it.\nBy morning, the thread is gone.",
    Inches(1), Inches(1.2), Inches(8), Inches(2.0),
    size=26, color=WHITE_DIM)

# Three breakdown pills
breakdown_items = ["Thoughts reset", "Good ideas die", "Continuity is lost"]
pill_y = Inches(3.8)
for i, item in enumerate(breakdown_items):
    px = Inches(1) + i * Inches(4.0)
    pill(sl, item, px, pill_y, Inches(3.6), Inches(0.65),
         fill=RGBColor(0x12,0x18,0x2c),
         border=RGBColor(0x28,0x38,0x66),
         text_color=WHITE, size=22)

txb(sl, "Every tool we use treats each day as a separate conversation.",
    Inches(1), Inches(5.0), Inches(9), Inches(0.5),
    size=16, color=WHITE_SUB)

# ── SLIDE 4 — Solution: Sakhi ─────────────────────────────────────────────────
sl = add_slide()

txb(sl, "Sakhi makes it all\ncome together.",
    Inches(0.8), Inches(1.2), Inches(11), Inches(3.0),
    size=60, bold=True, color=WHITE)

txb(sl, "You do not have to be driven by your thoughts.",
    Inches(0.8), Inches(4.6), Inches(10), Inches(0.5),
    size=22, color=WHITE_DIM)
txb(sl, "You can shape them.",
    Inches(0.8), Inches(5.15), Inches(10), Inches(0.5),
    size=22, bold=True, color=WHITE)

# ── SLIDE 5 — Stop Reacting. Start Shaping. ───────────────────────────────────
sl = add_slide()

txb(sl, "Stop reacting.", Inches(0.8), Inches(0.8), Inches(11), Inches(1.3),
    size=54, bold=True, color=WHITE)
txb(sl, "Start shaping.",  Inches(0.8), Inches(2.0), Inches(11), Inches(1.3),
    size=54, bold=True, color=WHITE_SUB)

txb(sl, "Sakhi learns how you think and evolves with you,\nso your decisions become clearer over time.",
    Inches(0.8), Inches(3.5), Inches(9), Inches(1.0),
    size=18, color=WHITE_DIM)

# Three pillars
pillars = [
    ("Pillar 1", "Builds a living model of you.", "It evolves over time."),
    ("Pillar 2", "Makes your life visible.", "Across time, as a whole."),
    ("Pillar 3", "Helps you act decisively.", "It learns from what follows."),
]
for i, (kicker, title, body) in enumerate(pillars):
    px = Inches(0.8) + i * Inches(4.15)
    py = Inches(4.7)
    rect(sl, px, py, Inches(3.9), Inches(2.2),
         fill_color=RGBColor(0x0c,0x14,0x28),
         line_color=RGBColor(0x22,0x32,0x58))
    txb(sl, kicker, px + Inches(0.15), py + Inches(0.15), Inches(3.6), Inches(0.25),
        size=9, bold=True, color=ACCENT)
    txb(sl, title,  px + Inches(0.15), py + Inches(0.45), Inches(3.6), Inches(0.5),
        size=16, bold=True, color=WHITE)
    txb(sl, body,   px + Inches(0.15), py + Inches(0.95), Inches(3.6), Inches(0.9),
        size=13, color=WHITE_DIM)

# ── SLIDE 6 — Life Occupancy ──────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "CONTINUITY  ·  VIDHYA · PROFILE", Inches(1), Inches(0.5))
txb(sl, "Life Occupancy", Inches(1), Inches(0.9), Inches(9), Inches(0.8),
    size=36, bold=True, color=WHITE)
txb(sl, "Bubble size reflects how much each thread has occupied your attention.",
    Inches(1), Inches(1.65), Inches(8), Inches(0.4),
    size=14, color=WHITE_DIM)

# Stat badge
rect(sl, Inches(10.5), Inches(0.8), Inches(2.0), Inches(0.9),
     fill_color=RGBColor(0x0c,0x14,0x28), line_color=RGBColor(0x22,0x32,0x58))
txb(sl, "5 threads", Inches(10.5), Inches(0.88), Inches(2.0), Inches(0.4),
    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Active", Inches(10.5), Inches(1.28), Inches(2.0), Inches(0.25),
    size=9, bold=True, color=WHITE_SUB, align=PP_ALIGN.CENTER)

# Bubbles — approximate positions
bubbles = [
    ("Start up",   "65%", "28 moments", Inches(4.0),  Inches(3.8),  Inches(2.5), Inches(2.5), True),
    ("Family",     "14%", "6 moments",  Inches(8.2),  Inches(2.0),  Inches(1.4), Inches(1.4), False),
    ("Caregiving", "7%",  "3 moments",  Inches(6.5),  Inches(5.2),  Inches(1.1), Inches(1.1), False),
    ("Career",     "7%",  "3 moments",  Inches(8.2),  Inches(4.4),  Inches(1.1), Inches(1.1), False),
    ("Self Care",  "7%",  "3 moments",  Inches(10.0), Inches(3.3),  Inches(1.1), Inches(1.1), False),
]
for label, share, moments, bx, by, bw, bh, accent in bubbles:
    sh = sl.shapes.add_shape(9, bx - bw/2, by - bh/2, bw, bh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x18,0x28,0x50) if accent else RGBColor(0x10,0x18,0x30)
    sh.line.color.rgb = ACCENT if accent else RGBColor(0x28,0x38,0x60)
    sh.line.width = Pt(1.2 if accent else 0.6)

    inner_left  = bx - bw/2 + Inches(0.05)
    inner_top   = by - bh/2 + Inches(0.05)
    inner_w     = bw - Inches(0.1)
    inner_h     = bh - Inches(0.1)
    name_size   = 12 if accent else 10
    share_size  = 28 if accent else 18
    sub_size    = 10 if accent else 8

    txb(sl, label, inner_left, inner_top + inner_h * 0.15, inner_w, Inches(0.3),
        size=name_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, share, inner_left, inner_top + inner_h * 0.38, inner_w, Inches(0.6),
        size=share_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, moments, inner_left, inner_top + inner_h * 0.72, inner_w, Inches(0.25),
        size=sub_size, color=WHITE_DIM, align=PP_ALIGN.CENTER)

# ── SLIDE 7 — Continuity Arc ──────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "START UP  ·  CONTINUITY ARC", Inches(0.6), Inches(0.3))
txb(sl, "104 days of continuity — the startup thread, mapped.",
    Inches(0.6), Inches(0.65), Inches(11), Inches(0.45),
    size=18, bold=True, color=WHITE)

arc_stages = [
    ("Where It Began", "An idea with ache",    "Day 12",  True),
    ("Phase 1",        "Founder question",      "Day 16",  False),
    ("Phase 2",        "Rhythm clicked",        "Day 19",  False),
    ("Phase 3",        "RAG felt brittle",      "Day 24",  False),
    ("Phase 4",        "Knowledge graph",       "Day 28",  False),
    ("Phase 5",        "Still uneasy",          "Day 35",  False),
    ("Phase 6",        "Clarity surfaced",      "Day 41",  False),
    ("Phase 7",        "Validation pressure",   "Day 58",  False),
    ("Phase 8",        "Build, not pitch",      "Day 66",  False),
    ("Phase 9",        "Continuity core",       "Day 72",  False),
    ("Phase 10",       "Product shape",         "Day 115", False),
    ("Where It Is Now","A visible direction",   "104 days",True),
]

cols = 4
rows = 3
card_w = Inches(2.9)
card_h = Inches(1.7)
start_x = Inches(0.6)
start_y = Inches(1.25)
gap_x = Inches(0.28)
gap_y = Inches(0.22)

for i, (label, title, meta, major) in enumerate(arc_stages):
    row = i // cols
    col_in_row = i % cols
    col = col_in_row if row % 2 == 0 else (cols - 1 - col_in_row)
    cx = start_x + col * (card_w + gap_x)
    cy = start_y + row * (card_h + gap_y)

    fill = RGBColor(0x0a,0x12,0x26) if major else RGBColor(0x08,0x0e,0x1c)
    border = RGBColor(0x28,0x48,0x80) if major else RGBColor(0x18,0x26,0x42)
    rect(sl, cx, cy, card_w, card_h, fill_color=fill, line_color=border,
         line_width_pt=1.0 if major else 0.5)

    txb(sl, label, cx + Inches(0.12), cy + Inches(0.1), card_w - Inches(0.24), Inches(0.2),
        size=7, bold=True, color=ACCENT if major else WHITE_SUB)
    txb(sl, title, cx + Inches(0.12), cy + Inches(0.32), card_w - Inches(0.24), Inches(0.4),
        size=13 if major else 11, bold=True, color=WHITE)
    txb(sl, meta,  cx + Inches(0.12), cy + Inches(0.8), card_w - Inches(0.24), Inches(0.2),
        size=8, bold=True, color=ACCENT)

# ── SLIDE 8 — Your life, connected across time ────────────────────────────────
sl = add_slide()

label_kicker(sl, "CONTINUITY", Inches(1), Inches(0.7))

txb(sl, "Your life, connected\nacross time.",
    Inches(1), Inches(1.1), Inches(9), Inches(2.8),
    size=54, bold=True, color=WHITE)

lines = [
    "Noise becomes pattern.",
    "Pattern becomes readable occupancy.",
    "The system gathers scattered moments and returns which threads have actually occupied your life.",
    "So it shows not just what is happening now, but what has been taking up space across time.",
    "This is where scattered life becomes coherent.",
]
for i, line in enumerate(lines):
    c = WHITE if i == len(lines) - 1 else WHITE_DIM
    bold = i == len(lines) - 1
    txb(sl, line, Inches(1), Inches(4.0) + i * Inches(0.5), Inches(10), Inches(0.45),
        size=15, color=c, bold=bold)

# ── SLIDE 9 — First Expression ────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "FIRST EXPRESSION OF SAKHI", Inches(1), Inches(0.5))

txb(sl, "The product starts where life\nactually happens.",
    Inches(1), Inches(0.9), Inches(9), Inches(1.6),
    size=38, bold=True, color=WHITE)

txb(sl, "Conversation, reflection, moments, and privacy form the first expression of Sakhi.",
    Inches(1), Inches(2.6), Inches(11), Inches(0.5),
    size=15, color=WHITE_DIM)

panels = [
    ("Conversation", "Start talking. Sakhi does the rest.", "It keeps the thread. It builds over time. It evolves with you."),
    ("Reflection",   "Your life becomes visible.",           "Across time and topics, everything connects. Patterns emerge from what you've lived."),
    ("Continuity",   "Your story starts to take shape.",    "Each thread becomes something you can return to, and build on."),
    ("Privacy",      "What you share stays yours.",         "End-to-end encrypted. No one reads your conversations. Not even Sakhi."),
]
for i, (label, headline, body) in enumerate(panels):
    px = Inches(0.6) + i * Inches(3.15)
    py = Inches(3.25)
    rect(sl, px, py, Inches(2.9), Inches(3.8),
         fill_color=RGBColor(0x0c,0x14,0x28),
         line_color=RGBColor(0x22,0x32,0x58))
    txb(sl, label, px + Inches(0.12), py + Inches(0.12), Inches(2.66), Inches(0.25),
        size=9, bold=True, color=ACCENT)
    txb(sl, headline, px + Inches(0.12), py + Inches(0.42), Inches(2.66), Inches(0.8),
        size=14, bold=True, color=WHITE)
    txb(sl, body, px + Inches(0.12), py + Inches(1.3), Inches(2.66), Inches(2.0),
        size=12, color=WHITE_DIM)

# ── SLIDE 10 — Go-to-Market ───────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "GO TO MARKET", Inches(1), Inches(0.5))
txb(sl, "Built for people\nmanaging complex lives.",
    Inches(1), Inches(0.9), Inches(8), Inches(1.8),
    size=38, bold=True, color=WHITE)

txb(sl, "Urban professionals, founders, and knowledge workers who carry more threads than any tool can hold.",
    Inches(1), Inches(2.85), Inches(9.5), Inches(0.5), size=15, color=WHITE_DIM)

gtm_cols = [
    ("WHO",       "The person who holds it all",
     "Executives, founders, caregivers — navigating career, family, health simultaneously.\nUrban India · 28–50."),
    ("ENTRY",     "Personal AI companion",
     "Mobile-first, voice-native. One conversation → thread → pattern → clarity."),
    ("EXPANSION", "Platform and network",
     "Enterprise continuity. Sakhi Mesh (shared context). API layer for developers."),
]
for i, (kicker, title, body) in enumerate(gtm_cols):
    px = Inches(0.8) + i * Inches(4.15)
    py = Inches(3.5)
    rect(sl, px, py, Inches(3.9), Inches(2.4),
         fill_color=RGBColor(0x0c,0x14,0x28),
         line_color=RGBColor(0x22,0x32,0x58))
    txb(sl, kicker, px + Inches(0.15), py + Inches(0.12), Inches(3.6), Inches(0.22),
        size=9, bold=True, color=ACCENT)
    txb(sl, title,  px + Inches(0.15), py + Inches(0.38), Inches(3.6), Inches(0.5),
        size=15, bold=True, color=WHITE)
    txb(sl, body,   px + Inches(0.15), py + Inches(0.92), Inches(3.6), Inches(1.2),
        size=12, color=WHITE_DIM)

# Tiers
tiers = [
    ("Free",       "30 conversations/month", "Try it, feel the difference",        False),
    ("Pro",        "₹999 / month",           "Unlimited memory, voice, reflection", True),
    ("Enterprise", "Custom",                 "Teams, shared context, API access",   False),
]
tier_y = Inches(6.1)
for i, (name, price, tag, highlight) in enumerate(tiers):
    tx = Inches(0.8) + i * Inches(4.15)
    fill = RGBColor(0x0e,0x1a,0x3a) if highlight else RGBColor(0x0a,0x12,0x22)
    border = RGBColor(0x38,0x62,0xb0) if highlight else RGBColor(0x20,0x30,0x50)
    rect(sl, tx, tier_y, Inches(3.9), Inches(1.1), fill_color=fill, line_color=border)
    txb(sl, name,  tx + Inches(0.15), tier_y + Inches(0.08), Inches(3.6), Inches(0.3),
        size=11, bold=True, color=ACCENT if highlight else WHITE_DIM)
    txb(sl, price, tx + Inches(0.15), tier_y + Inches(0.38), Inches(3.6), Inches(0.35),
        size=16, bold=True, color=WHITE)
    txb(sl, tag,   tx + Inches(0.15), tier_y + Inches(0.73), Inches(3.6), Inches(0.28),
        size=11, color=WHITE_SUB)

# ── SLIDE 11 — Traction Target ────────────────────────────────────────────────
sl = add_slide()

txb(sl, "10,000", Inches(0.8), Inches(1.2), Inches(9), Inches(2.8),
    size=110, bold=True, color=WHITE)
txb(sl, "active users — Year 1 target.",
    Inches(0.8), Inches(4.1), Inches(9), Inches(0.6),
    size=26, bold=True, color=WHITE)
txb(sl, "100K is the Series A signal.",
    Inches(0.8), Inches(4.85), Inches(9), Inches(0.5),
    size=18, color=WHITE_DIM)
txb(sl, "15% D30 retention at scale proves the loop.",
    Inches(0.8), Inches(5.35), Inches(9), Inches(0.5),
    size=18, color=WHITE_DIM)

# ── SLIDE 12 — Vision ─────────────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "VISION", Inches(1), Inches(0.5))
txb(sl, "Sakhi makes your life seen,\nunderstood, and actionable.",
    Inches(1), Inches(0.9), Inches(9), Inches(2.0),
    size=40, bold=True, color=WHITE)

voice_pillars = [
    ("I see you.",         "Build a living model of you evolving over time."),
    ("I understand you.",  "Make your life visible across time, as a whole."),
    ("I act for you.",     "Help you act decisively and learn from what follows."),
]
for i, (voice, desc) in enumerate(voice_pillars):
    px = Inches(0.8) + i * Inches(4.15)
    py = Inches(3.2)
    rect(sl, px, py, Inches(3.9), Inches(2.8),
         fill_color=RGBColor(0x0c,0x14,0x28),
         line_color=RGBColor(0x22,0x32,0x58))
    txb(sl, voice, px + Inches(0.15), py + Inches(0.2), Inches(3.6), Inches(0.9),
        size=26, bold=True, color=WHITE)
    txb(sl, desc,  px + Inches(0.15), py + Inches(1.2), Inches(3.6), Inches(1.2),
        size=14, color=WHITE_DIM)

txb(sl, "— Sakhi", Inches(1), Inches(6.2), Inches(5), Inches(0.35),
    size=12, bold=True, color=ACCENT)

# ── SLIDE 13 — Founders Bridge ────────────────────────────────────────────────
sl = add_slide()

txb(sl, "The Minds\nBehind...",
    Inches(1), Inches(2.0), Inches(11), Inches(3.0),
    size=70, bold=True, color=WHITE)

label_kicker(sl, "FOUNDERS", Inches(1), Inches(5.2))

# ── SLIDE 14 — Vidhya ────────────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "CO-FOUNDER & CEO", Inches(1), Inches(0.4))
txb(sl, "Vidhya", Inches(1), Inches(0.75), Inches(6), Inches(0.7),
    size=32, bold=True, color=WHITE)
txb(sl, "Built systems for companies. Now building one for humans.",
    Inches(1), Inches(1.45), Inches(6.5), Inches(0.4),
    size=14, color=WHITE_DIM)

# Opening quote
rect(sl, Inches(1), Inches(2.0), Inches(6.8), Inches(1.4),
     fill_color=RGBColor(0x0c,0x14,0x2c), line_color=RGBColor(0x28,0x42,0x78))
txb(sl, '"I\'ve spent 20+ years helping organizations make better decisions.\nI realized we haven\'t solved this for individuals."',
    Inches(1.15), Inches(2.1), Inches(6.5), Inches(1.2),
    size=15, color=WHITE, italic=True)

# Arc steps
vidhya_arc = [
    ("Operator at Scale",        "Worked alongside CEOs and COOs, building systems that turned ambiguity into structured decisions."),
    ("Personal Inflection Point","In 2024, caregiving, leadership, and life complexity collided. What was missing was continuity at the level of a real human life."),
    ("Insight to Sakhi",         "Small, personalized interventions changed everything. Timing and personalization mattered more than generic advice."),
]
for i, (label, body) in enumerate(vidhya_arc):
    px = Inches(1) + i * Inches(4.1)
    py = Inches(3.65)
    rect(sl, px, py, Inches(3.85), Inches(2.3),
         fill_color=RGBColor(0x0a,0x12,0x24), line_color=RGBColor(0x20,0x30,0x5c))
    txb(sl, label, px + Inches(0.12), py + Inches(0.12), Inches(3.6), Inches(0.4),
        size=12, bold=True, color=WHITE)
    txb(sl, body,  px + Inches(0.12), py + Inches(0.55), Inches(3.6), Inches(1.5),
        size=12, color=WHITE_DIM)

txb(sl, '"Sakhi is the system I wish existed when I needed it most."',
    Inches(1), Inches(6.2), Inches(11), Inches(0.45),
    size=14, bold=True, color=WHITE, italic=True)

# ── SLIDE 15 — Ravi ───────────────────────────────────────────────────────────
sl = add_slide()

label_kicker(sl, "CO-FOUNDER & CTO", Inches(1), Inches(0.4))
txb(sl, "Ravi Shankar", Inches(1), Inches(0.75), Inches(7), Inches(0.7),
    size=32, bold=True, color=WHITE)
txb(sl, "Built systems across engineering, product, and AI. Now building one for the self.",
    Inches(1), Inches(1.45), Inches(7), Inches(0.4), size=14, color=WHITE_DIM)

rect(sl, Inches(1), Inches(2.0), Inches(6.8), Inches(1.2),
     fill_color=RGBColor(0x0c,0x14,0x2c), line_color=RGBColor(0x28,0x42,0x78))
txb(sl, '"I\'m a systems thinker at heart, grounded in deep technical expertise\nand driven to simplify complexity."',
    Inches(1.15), Inches(2.1), Inches(6.5), Inches(1.0),
    size=15, color=WHITE, italic=True)

ravi_arc = [
    ("Evolution",    "Started with engineering; kept moving closer to the question of what actually makes systems work."),
    ("Realization",  "Systems succeed because people understand, trust, and use them — not just because they are built well."),
    ("Expansion",    "Moved across engineering, product, and product marketing. Yoga and meditation deepened how he thinks about human behavior."),
    ("Convergence",  "Technical depth, systems thinking, product narrative, and lived understanding converge in building Sakhi."),
]
for i, (label, body) in enumerate(ravi_arc):
    px = Inches(0.8) + i * Inches(3.12)
    py = Inches(3.5)
    rect(sl, px, py, Inches(2.95), Inches(2.5),
         fill_color=RGBColor(0x0a,0x12,0x24), line_color=RGBColor(0x20,0x30,0x5c))
    txb(sl, label, px + Inches(0.12), py + Inches(0.12), Inches(2.7), Inches(0.35),
        size=12, bold=True, color=WHITE)
    txb(sl, body,  px + Inches(0.12), py + Inches(0.5), Inches(2.7), Inches(1.7),
        size=12, color=WHITE_DIM)

txb(sl, '"That gives me the clarity to build Sakhi."',
    Inches(1), Inches(6.2), Inches(11), Inches(0.45),
    size=14, bold=True, color=WHITE, italic=True)

# ── SLIDE 16 — The Ask ────────────────────────────────────────────────────────
sl = add_slide()

txb(sl, "This is just\nthe beginning.",
    Inches(1), Inches(0.9), Inches(11), Inches(2.6),
    size=60, bold=True, color=WHITE)

txb(sl, "If this resonates, let's build this together.\nWe share everything: vision, GTM, early product, and roadmap.",
    Inches(1), Inches(3.7), Inches(10), Inches(0.8),
    size=17, color=WHITE_DIM)

# Two CTA cards
cta = [
    ("Collaborate", "If you believe we should shape our lives, not just react to them, come build this with us."),
    ("Invest",      "For investors who see this space the way we do, we share everything: vision, GTM, early product, and roadmap."),
]
for i, (title, body) in enumerate(cta):
    cx = Inches(1) + i * Inches(5.8)
    cy = Inches(4.7)
    rect(sl, cx, cy, Inches(5.4), Inches(1.8),
         fill_color=RGBColor(0x0c,0x14,0x28), line_color=RGBColor(0x22,0x34,0x60))
    txb(sl, title, cx + Inches(0.18), cy + Inches(0.15), Inches(5.0), Inches(0.35),
        size=16, bold=True, color=WHITE)
    txb(sl, body,  cx + Inches(0.18), cy + Inches(0.55), Inches(5.0), Inches(1.0),
        size=12, color=WHITE_DIM)

txb(sl, "For collaboration or investment:",
    Inches(1), Inches(6.6), Inches(8), Inches(0.3),
    size=10, bold=True, color=WHITE_SUB)
txb(sl, "founders@sakhiintelligence.com",
    Inches(1), Inches(6.92), Inches(8), Inches(0.38),
    size=16, bold=True, color=ACCENT)
txb(sl, "Sakhi · 2026",
    Inches(10.5), Inches(7.05), Inches(2.0), Inches(0.3),
    size=9, color=WHITE_SUB, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(
    os.path.dirname(os.path.dirname(__file__)),
    "docs", "Sakhi_Pitch_Deck.pptx"
)
prs.save(out_path)
print(f"Saved → {out_path}")
print(f"Slides: {len(prs.slides)}")
